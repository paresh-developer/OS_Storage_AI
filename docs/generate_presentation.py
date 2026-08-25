"""Generates docs/presentation.pptx -- an academic-style overview deck for
the whole project, using python-pptx (a light, pure-Python dependency
needed only for this script -- `.venv/bin/pip install python-pptx`, not
part of the app's own requirements).

Run from the repo root: .venv/bin/python docs/generate_presentation.py

The screenshots under docs/screenshots/ are real captures from a genuine
scan/run (see each slide's footer), not mock-ups -- if the GUI changes
enough to make them stale, recapture them (offscreen Qt + QWidget.grab(),
same approach used to build them originally) rather than hand-editing
these images. Re-run this script after any such recapture, and after any
structural change already reflected in docs/architecture.png."""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SCREENSHOTS = os.path.join(REPO, "docs", "screenshots")
OUT_PATH = os.path.join(REPO, "docs", "presentation.pptx")

DARK = RGBColor(0x1F, 0x2D, 0x4E)
ACCENT = RGBColor(0x33, 0x53, 0x8A)
LIGHT = RGBColor(0xDB, 0xE4, 0xF0)
BG = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x55, 0x5F, 0x70)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)


def new_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # fully blank layout


def add_rect(slide, left, top, width, height, color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line:
        shape.line.color.rgb = color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_textbox(slide, left, top, width, height, text, size=18, color=DARK, bold=False,
                 align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return box


def add_header(slide, kicker, title):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.14), ACCENT)
    add_textbox(slide, MARGIN, Inches(0.35), SLIDE_W - 2 * MARGIN, Inches(0.4),
                kicker, size=14, color=ACCENT, bold=True)
    # A long title wraps to two lines at this width/size -- give it the extra
    # room rather than letting it overlap the content below (measured
    # empirically: titles over ~48 chars wrapped in the LibreOffice-rendered
    # preview at 30pt bold on this slide width).
    two_line = len(title) > 48
    title_top = Inches(0.72)
    title_h = Inches(1.2) if two_line else Inches(0.7)
    add_textbox(slide, MARGIN, title_top, SLIDE_W - 2 * MARGIN, title_h,
                title, size=30, color=DARK, bold=True)
    return title_top + title_h + Inches(0.2)


def add_footer(slide, text):
    add_textbox(slide, MARGIN, SLIDE_H - Inches(0.5), SLIDE_W - 2 * MARGIN, Inches(0.35),
                text, size=11, color=GREY, italic=True)


def add_bullets(slide, left, top, width, height, items, size=17, gap_after=10, color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(gap_after)
        run = p.add_run()
        marker = "•  " if level == 0 else "-  "
        run.text = marker + text
        run.font.size = Pt(size if level == 0 else size - 2)
        run.font.color.rgb = color if level == 0 else GREY
        run.font.name = "Calibri"
    return box


def fit_picture(slide, path, left, top, max_w, max_h, border=False, card=False, card_pad=Inches(0.15)):
    with Image.open(path) as im:
        w_px, h_px = im.size
    ratio = w_px / h_px
    box_ratio = max_w / max_h
    if ratio > box_ratio:
        w = max_w
        h = Emu(int(max_w / ratio))
    else:
        h = max_h
        w = Emu(int(max_h * ratio))
    x = left + Emu(int((max_w - w) / 2))
    y = top + Emu(int((max_h - h) / 2))
    if card:
        add_rect(
            slide,
            x - card_pad, y - card_pad, w + 2 * card_pad, h + 2 * card_pad,
            RGBColor(0xF2, 0xF4, 0xF8),
        )
    pic = slide.shapes.add_picture(path, x, y, width=w, height=h)
    if border:
        pic.line.color.rgb = RGBColor(0xAA, 0xB4, 0xC8)
        pic.line.width = Pt(1)
    return x, y, w, h


def title_slide(prs, title, subtitle, footer):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK)
    add_rect(slide, 0, Inches(4.55), SLIDE_W, Inches(0.06), ACCENT)
    add_textbox(slide, Inches(1), Inches(2.7), SLIDE_W - Inches(2), Inches(1.3),
                title, size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.75), SLIDE_W - Inches(2), Inches(0.8),
                subtitle, size=20, color=LIGHT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), SLIDE_H - Inches(0.9), SLIDE_W - Inches(2), Inches(0.5),
                footer, size=13, color=LIGHT, align=PP_ALIGN.CENTER, italic=True)
    return slide


def bullets_slide(prs, kicker, title, items, footer=None, size=18):
    slide = blank_slide(prs)
    top = add_header(slide, kicker, title)
    add_bullets(slide, MARGIN, top, SLIDE_W - 2 * MARGIN, SLIDE_H - top - Inches(0.7), items, size=size)
    if footer:
        add_footer(slide, footer)
    return slide


def split_slide(prs, kicker, title, items, image_path, footer=None, size=16.5):
    slide = blank_slide(prs)
    top = add_header(slide, kicker, title)
    left_w = Inches(6.3)
    add_bullets(slide, MARGIN, top, left_w, SLIDE_H - top - Inches(0.7), items, size=size)
    img_left = MARGIN + left_w + Inches(0.4)
    img_w = SLIDE_W - img_left - MARGIN
    fit_picture(slide, image_path, img_left, top + Inches(0.3), img_w, SLIDE_H - top - Inches(1.3),
                border=True, card=True)
    if footer:
        add_footer(slide, footer)
    return slide


def image_slide(prs, kicker, title, image_path, footer=None):
    slide = blank_slide(prs)
    top = add_header(slide, kicker, title)
    fit_picture(slide, image_path, MARGIN, top, SLIDE_W - 2 * MARGIN, SLIDE_H - top - Inches(0.7))
    if footer:
        add_footer(slide, footer)
    return slide


def section_slide(prs, number, title, subtitle):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, ACCENT)
    add_textbox(slide, Inches(1), Inches(2.6), Inches(2), Inches(1.4), number,
                size=64, color=RGBColor(0x9F, 0xB4, 0xD8), bold=True)
    add_textbox(slide, Inches(1), Inches(3.6), SLIDE_W - Inches(2), Inches(1.0),
                title, size=36, color=WHITE, bold=True)
    add_textbox(slide, Inches(1), Inches(4.45), SLIDE_W - Inches(2), Inches(0.8),
                subtitle, size=17, color=LIGHT)
    return slide


def main():
    prs = new_deck()

    title_slide(
        prs,
        "Storage AI",
        "An Intelligent, Explainable Storage Cleanup Assistant",
        "Academic project — AI at the Application Level  ·  storage_ai/",
    )

    bullets_slide(
        prs, "MOTIVATION", "The problem this project solves",
        [
            "Storage fills up silently: duplicate files, stale downloads, unbounded logs and caches.",
            "Manual triage is slow and risky — which files are truly unused? which app owns this data?",
            "Goal: an assistant that finds, explains, and safely acts on cleanup opportunities.",
            "Built as an academic project on “AI at the Application Level.”",
        ],
        footer="Storage AI  ·  Motivation",
    )

    bullets_slide(
        prs, "DESIGN PHILOSOPHY", "Classical, explainable AI — not an LLM by default",
        [
            "Every recommendation traces back to a concrete, inspectable reason.",
            "Content hashing, weakly-supervised classification, linear regression, K-means clustering.",
            "A local LLM appears exactly once, as a narrowly-scoped last resort (Section 8) — never the default path.",
            "No API keys, no GPU, no network calls at runtime — fully offline and reproducible.",
        ],
        footer="Storage AI  ·  Design philosophy",
    )

    image_slide(
        prs, "ARCHITECTURE", "Two orchestrators, one shared store",
        os.path.join(REPO, "docs", "architecture.png"),
        footer="Storage AI  ·  Architecture  ·  see docs/architecture.png",
    )

    bullets_slide(
        prs, "CORE ANALYSIS", "Duplicate detection & unused-file scoring",
        [
            "Duplicate detection: size → partial-hash → full-hash funnel — only real candidates get fully read.",
            "Unused-file scoring: a RandomForestClassifier trained on weakly-labeled access patterns.",
            "Falls back to a pure heuristic when there isn't enough data to train a real model.",
            "Produces a 0–1 “likely unused” score per file — informational only, never auto-deleted.",
        ],
        footer="Storage AI  ·  Core analysis",
    )

    bullets_slide(
        prs, "CORE ANALYSIS", "Growth forecasting & file clustering",
        [
            "Storage growth forecasting: linear regression over real scan history, or a file-timestamp pseudo-history on a first scan.",
            "K-means clustering groups files by size and staleness into archetypes like “Large & Stale.”",
            "An unsupervised cross-check, independent of the unused-file classifier — not a duplicate of it.",
            "“No growth trend detected” is a correct answer, not a bug, when there's no signal to extrapolate.",
        ],
        footer="Storage AI  ·  Core analysis",
    )

    split_slide(
        prs, "CORE ANALYSIS", "Path classification & ranked recommendations",
        [
            "Every file tagged: system · log · cache · application_data · user_data · trash · other.",
            "Known-service bonus label (e.g. /var/lib/postgresql → PostgreSQL) with tailored advice.",
            "Duplicates, unused files, storage warnings, and category advisories merge into one ranked list.",
            "Sorted by estimated space recovered, each row backed by a concrete reason.",
        ],
        os.path.join(SCREENSHOTS, "recommendations.png"),
        footer="Storage AI  ·  Real screenshot: a genuine demo scan (duplicate + 2 unused files, 1.9 MB recoverable)",
    )

    split_slide(
        prs, "CORE ANALYSIS", "Visualization — and provenance on demand",
        [
            "Dashboard, File Types, 30-day Forecast, Folder treemap, and Cluster scatter — one story per chart.",
            "Every chart's ℹ button now shows the real directories behind a total, not just an opaque number.",
            "E.g. “User data: 1.3 GB” becomes an actual, inspectable list of folders on click.",
            "legend_detail.py: one small module, reused unchanged across four structurally different charts.",
        ],
        os.path.join(SCREENSHOTS, "folders.png"),
        footer="Storage AI  ·  Real screenshot: top-level folder treemap from a genuine demo scan",
    )

    bullets_slide(
        prs, "SAFETY", "Safe by default, reversible by design",
        [
            "Cleanup actions send files to the OS trash or a local dated archive folder — never a hard delete.",
            "Files under a live service's data directory, or real OS system paths, are never offered as candidates.",
            "An advisory can only ever suggest — no advisory string can itself trigger a delete or archive.",
            "A storage tool that's occasionally wrong about “unused” must stay recoverable.",
        ],
        footer="Storage AI  ·  Safety",
    )

    bullets_slide(
        prs, "LIVE MONITORING", "Real-time activity monitoring (Section 6)",
        [
            "Cross-platform create/modify/delete watcher (watchdog) — runs from the GUI or as a standalone service.",
            "Rate-based trend alerts: large file added, rapid deletes by one user, activity bursts.",
            "Alerts fire on plain, explainable thresholds — not a learned model.",
            "Optional Linux auditd integration upgrades attribution from “who owns this file” to “who actually did this.”",
        ],
        footer="Storage AI  ·  Live activity monitoring",
    )

    bullets_slide(
        prs, "APPLICATION DISCOVERY", "How does it find an app it's never been taught? (Section 8)",
        [
            "The core question: config files live in different places per app — how to find storagePath generically?",
            "1. Live process introspection (/proc) — real command-line flags and open files, not a claim.",
            "2. Structured config-file parsing — JSON, YAML, TOML, INI, fuzzy-matched for path-like settings.",
            "3. The curated known-service table — classification, not discovery (the opposite direction).",
            "4. An optional local, CPU-only extractive-QA model — a last resort, only for the residual case.",
        ],
        footer="Storage AI  ·  Application discovery  ·  docs/METHODOLOGY.md Section 8",
        size=17,
    )

    bullets_slide(
        prs, "APPLICATION DISCOVERY", "Real bugs found — not hypothesized",
        [
            "Question-phrasing brittleness: identical text, different question wording, wildly different extraction quality.",
            "Span imprecision: naive token decoding fragmented paths (“mongo - storage”) — fixed with offset-mapping.",
            ("A verified hallucination: the model reported 0.47 confidence extracting an answer from text with NO path in it at all.", 1),
            "Fix: a mandatory, deterministic path-shape check — independent of confidence — locked in as a regression test.",
        ],
        footer="Storage AI  ·  Application discovery  ·  empirically calibrated, not just implemented",
    )

    split_slide(
        prs, "APPLICATION DISCOVERY", "App Data Suggestions — batched across the whole machine",
        [
            "app_suggestions.py reuses app_discovery.py unchanged, once per distinct running process.",
            "Only surfaces a finding if it maps to real advisory text — stays a short, actionable list.",
            "The discovered path's real on-disk usage is measured and flagged if unusually large.",
            "Independent Run / Rerun and Stop controls, cancellable mid-run — same pattern as the main scan.",
        ],
        os.path.join(SCREENSHOTS, "app_suggestions.png"),
        footer="Storage AI  ·  App Data Suggestions tab  ·  real classification/advice/severity logic",
    )

    bullets_slide(
        prs, "APPLICATION DISCOVERY", "Worked example: a 1 TB MongoDB data directory",
        [
            "mongod is found running with --dbpath /var/lib/mongodb (tier 1, live process introspection).",
            "Classified via the curated table as MongoDB → generic advice: \"Consider MongoDB's built-in log rotation...\"",
            "A real filesystem walk of that path measures its actual size — deferred until after every cheaper filter already passed, so it's not wasted on findings nobody will see.",
            "1 TB crosses the 50 GB \"critical\" threshold → the row is flagged ❗, bold, colored, and sorted to the top, ahead of every other finding regardless of discovery confidence.",
            "Previously the same finding would show identical advice whether it used 10 MB or 10 TB — the size check is what turns a generic tip into an actual priority.",
        ],
        footer="Storage AI  ·  Worked example  ·  LARGE_SIZE_BYTES = 5 GB, CRITICAL_SIZE_BYTES = 50 GB",
    )

    bullets_slide(
        prs, "NOVELTY", "What's actually new here — for the record",
        [
            "None of the individual algorithms are new — hashing, random forest, K-means, linear regression, extractive QA are all textbook.",
            "Reliability-ordered discovery, not LLM-first — cheapest/most certain signal tried first.",
            "A deterministic gate around every non-deterministic component — applied independently in 3 unrelated subsystems.",
            "A single-app lookup reused unchanged as a whole-system inventory (zero new discovery logic).",
            "On-demand provenance for every displayed total, not just ML outputs.",
        ],
        footer="Storage AI  ·  Novelty  ·  docs/METHODOLOGY.md Section 9",
        size=17,
    )

    bullets_slide(
        prs, "HONESTY", "Known limitations",
        [
            "Duplicate detection is exact-match only — near-duplicates are out of scope.",
            "Known-service labels are default-location hints, not verified facts — a custom install won't match.",
            "Without --enable-audit, activity attribution is file-ownership based, not per-operation.",
            "The LLM tier's own confidence score is not trustworthy alone — verified, not assumed.",
            "No feedback loop yet records a user's accept/reject decision on a suggestion.",
        ],
        footer="Storage AI  ·  Known limitations  ·  docs/METHODOLOGY.md",
    )

    bullets_slide(
        prs, "VERIFICATION", "Tested for real, not just implemented",
        [
            "215 automated tests — non-GUI logic, plus GUI logic exercised via offscreen Qt.",
            "Real subprocesses tested for /proc-based discovery — not mocked.",
            "The LLM tier tested against the real model — its calibration bugs are permanent regression tests.",
            "Fully offline: no network calls at runtime, no GPU, no API keys, no external services.",
        ],
        footer="Storage AI  ·  Verification",
    )

    bullets_slide(
        prs, "WHAT'S NEXT", "Future work",
        [
            "Multi-folder & scheduled scanning; incremental scanning for large, slow-changing trees.",
            "Detect and prefer an already-running local LLM (e.g. Ollama) instead of always bundling one.",
            "A feedback-loop table recording accept/reject on every suggestion — the prerequisite for personalization.",
            "Multi-host aggregation — a fleet-wide storage inventory across several machines.",
        ],
        footer="Storage AI  ·  Future work  ·  see future_plan.md",
    )

    closing = blank_slide(prs)
    add_rect(closing, 0, 0, SLIDE_W, SLIDE_H, DARK)
    add_textbox(closing, Inches(1), Inches(2.9), SLIDE_W - Inches(2), Inches(1.0),
                "Thank you", size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(closing, Inches(1), Inches(3.9), SLIDE_W - Inches(2), Inches(0.6),
                "Questions?", size=20, color=LIGHT, align=PP_ALIGN.CENTER)
    add_textbox(closing, Inches(1), Inches(5.0), SLIDE_W - Inches(2), Inches(0.8),
                "docs/METHODOLOGY.md  ·  ReadMe.md  ·  future_plan.md",
                size=14, color=LIGHT, align=PP_ALIGN.CENTER, italic=True)

    prs.save(OUT_PATH)
    print("saved", OUT_PATH, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))


if __name__ == "__main__":
    main()
